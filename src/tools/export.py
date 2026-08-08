"""交付物导出：把生成的报告（Markdown）+ 指标，导出为 JSON / Markdown / PDF / PPTX。

- PDF 用 reportlab 内置 STSong-Light CID 字体，原生支持中文，无需额外字体文件；
- PPTX 用 python-pptx，按 `## ` 小节拆分幻灯片；
- 所有函数返回 (bytes, media_type, file_ext)，便于 Web 端点直接作为下载流。
"""
from __future__ import annotations

import json
import os
import re
from pathlib import Path
from typing import Dict, Tuple

# 图片目录：src/tools/export.py -> 项目根/outputs/images
IMAGES_ROOT = Path(__file__).resolve().parents[2] / "outputs" / "images"


def _resolve_image(url: str):
    """把 Markdown 图片 URL 解析为本地文件路径；找不到返回 None。"""
    if not url:
        return None
    name = url.rsplit("/", 1)[-1]
    candidates = []
    if url.startswith("/images/"):
        candidates.append(IMAGES_ROOT / name)
    candidates.append(Path(url))  # 绝对路径或相对路径
    for c in candidates:
        try:
            if c.exists():
                return str(c)
        except Exception:
            continue
    return None


def _inject_images(report_md: str) -> str:
    """扫描 images 目录，把封面图与章节配图以 Markdown 图片语法注入报告，
    供 PDF/PPTX 渲染时识别并插入。章节图按 `## N.` 编号匹配 sec-<N>.png。"""
    if not IMAGES_ROOT.exists():
        return report_md
    lines = report_md.splitlines()
    out = []
    if (IMAGES_ROOT / "cover.png").exists():
        out.append("![封面背景](/images/cover.png)")
    for ln in lines:
        out.append(ln)
        m = re.match(r"^##\s+(\d+)\.\s", ln)
        if m and (IMAGES_ROOT / f"sec-{m.group(1)}.png").exists():
            out.append(f"![第{m.group(1)}章配图](/images/sec-{m.group(1)}.png)")
    return "\n".join(out)


def _escape(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _inline_md(s: str) -> str:
    """转 **加粗** 为 <b>，并转义 HTML 特殊字符。"""
    s = _escape(s)
    return re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", s)


def export_report(report_md: str, metrics: Dict, fmt: str) -> Tuple[bytes, str, str]:
    fmt = (fmt or "md").lower()
    if fmt == "json":
        data = json.dumps({"report": report_md, "metrics": metrics}, ensure_ascii=False, indent=2)
        return data.encode("utf-8"), "application/json", "json"
    if fmt == "md":
        return report_md.encode("utf-8"), "text/markdown", "md"
    if fmt == "pdf":
        return _to_pdf(_inject_images(report_md), metrics), "application/pdf", "pdf"
    if fmt == "pptx":
        return _to_pptx(_inject_images(report_md), metrics), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx"
    raise ValueError(f"不支持的导出格式: {fmt}")


# ---------- PDF ----------
def _to_pdf(report_md: str, metrics: Dict) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont


    def _pdf_image(path, width=460):
        ir = ImageReader(path)
        iw, ih = ir.getSize()
        h = width * ih / float(iw) if iw else width
        return Image(path, width=width, height=h)

    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    FONT = "STSong-Light"

    styles = getSampleStyleSheet()
    title = ParagraphStyle("title", parent=styles["Title"], fontName=FONT, fontSize=18, leading=24)
    h1 = ParagraphStyle("h1", parent=styles["Heading1"], fontName=FONT, fontSize=14, leading=20, spaceBefore=10)
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], fontName=FONT, fontSize=12, leading=17, spaceBefore=8)
    body = ParagraphStyle("body", parent=styles["BodyText"], fontName=FONT, fontSize=10, leading=15)
    bullet = ParagraphStyle("bullet", parent=body, leftIndent=14, bulletIndent=4, spaceBefore=2)

    import io

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, title="研究报告")
    flow = []
    for line in report_md.splitlines():
        line = line.rstrip()
        if not line.strip():
            flow.append(Spacer(1, 4))
            continue
        # 图片行：![alt](url) -> 插入图片（封面大图 / 章节配图）
        if line.startswith("![") and "](" in line:
            m = re.match(r"^!\[(.*?)\]\((.*?)\)$", line.strip())
            if m:
                path = _resolve_image(m.group(2))
                if path:
                    try:
                        flow.append(_pdf_image(path))
                        flow.append(Spacer(1, 6))
                    except Exception:
                        pass
            continue
        if line.startswith("## "):
            flow.append(Paragraph(_inline_md(line[3:]), h2))
        elif line.startswith("# "):
            flow.append(Paragraph(_inline_md(line[2:]), title))
        elif re.match(r"^[-*] ", line):
            flow.append(Paragraph("• " + _inline_md(line[2:]), bullet))
        else:
            flow.append(Paragraph(_inline_md(line), body))

    # 指标页
    flow.append(Spacer(1, 12))
    flow.append(Paragraph("生成指标", h2))
    for k, v in (metrics or {}).items():
        if k in ("trace", "cost"):
            continue
        flow.append(Paragraph(f"• {k}: {v}", bullet))

    doc.build(flow)
    return buf.getvalue()


# ---------- PPTX ----------
def _to_pptx(report_md: str, metrics: Dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor

    prs = Presentation()
    lines = report_md.splitlines()

    # 标题取首个 # 标题
    title_text = "研究报告"
    for ln in lines:
        if ln.startswith("# "):
            title_text = ln[2:].strip()
            break
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "多智能体协作生成 · 带引用深度研究报告"
    # 封面背景图：作为标题幻灯片底层背景，标题文字改浅色以适配深色底图
    cover = _resolve_image("/images/cover.png")
    if cover:
        try:
            pic = slide.shapes.add_picture(cover, 0, 0, width=prs.slide_width, height=prs.slide_height)
            sp_tree = slide.shapes._spTree
            sp_tree.remove(pic._element)
            sp_tree.insert(2, pic._element)  # 移到最底层，避免盖住标题
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xDD, 0xE6, 0xF5)
        except Exception:
            pass

    # 按 ## 拆分内容幻灯片
    sections = []
    cur_title, cur_body = None, []
    for ln in lines:
        if ln.startswith("## "):
            if cur_title is not None:
                sections.append((cur_title, cur_body))
            cur_title = ln[3:].strip()
            cur_body = []
        elif ln.strip():
            cur_body.append(ln.strip())
    if cur_title is not None:
        sections.append((cur_title, cur_body))

    if not sections:
        sections = [("报告内容", [ln.strip() for ln in lines if ln.strip()])]

    for sec_title, sec_body in sections:
        # 长章节自动分页，避免单页溢出被裁切
        lines = [b for b in sec_body
                 if re.sub(r"^[-*] ", "", re.sub(r"\*\*", "", re.sub(r"^#+\s*", "", b))).strip()]
        chunk_size = 9
        chunks = [lines[i:i + chunk_size] for i in range(0, max(len(lines), 1), chunk_size)] or [[""]]
        for ci, chunk in enumerate(chunks):
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = sec_title + (f"（{ci + 1}/{len(chunks)}）" if len(chunks) > 1 else "")
            tf = s.placeholders[1].text_frame
            tf.word_wrap = True
            first = True
            for b in chunk:
                b = re.sub(r"^[-*] ", "", b)
                b = re.sub(r"\*\*", "", b)
                b = re.sub(r"^#+\s*", "", b)
                if not b.strip():
                    continue
                # 章节配图行：![alt](url) -> 在幻灯片底部插入图片
                if b.strip().startswith("![") and "](" in b:
                    m = re.match(r"^!\[(.*?)\]\((.*?)\)$", b.strip())
                    if m:
                        path = _resolve_image(m.group(2))
                        if path:
                            try:
                                slide.shapes.add_picture(path, Inches(0.4), Inches(5.2), width=Inches(6.7))
                            except Exception:
                                pass
                    continue
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.text = b
                p.font.size = Pt(16)
                first = False

    # 指标幻灯片
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "生成指标"
    tf = s.placeholders[1].text_frame
    tf.word_wrap = True
    first = True
    for k, v in (metrics or {}).items():
        if k in ("trace", "cost"):
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = f"{k}: {v}"
        p.font.size = Pt(16)
        first = False

    import io
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
